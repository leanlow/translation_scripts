#! /usr/bin/env python
import csv, sys, numpy, math, pandas, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path

prs = Presentation()
prs.slide_height = Inches(5.63)
prs.slide_width = Inches(10)
blank_slide_layout = prs.slide_layouts[6]

width = Inches(7.28)
height = Inches(0.52)

textfile1 = open(sys.argv[1]).read().splitlines()

count = 0

for line in textfile1:
	slide = prs.slides.add_slide(blank_slide_layout)
	#(left, top, width, height)
	txBox = slide.shapes.add_textbox(Inches(2.63), Inches(2.3), width, height)
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.alignment = PP_ALIGN.CENTER
	run = p.add_run()
	run.text = line
	font = run.font
	font.name = 'Throw My Hands Up in the Air Regular'
	font.size = Pt(35)

prs.save('/Users/loser/Downloads/finished_slides.pptx')

