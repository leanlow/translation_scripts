#! /usr/bin/env python
import csv, sys, numpy, math, pandas, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path

prs = Presentation()
prs.slide_height = Inches(5.63)
prs.slide_width = Inches(10)
blank_slide_layout = prs.slide_layouts[6]

width = Inches(10)
#(left, top, width, height)
textfile1 = open(sys.argv[1]).read().splitlines()
textfile2 = open(sys.argv[2]).read().splitlines()
textfile3 = open(sys.argv[3]).read().splitlines()

count = 0
count1 = 0

for line in textfile1:
	slide = prs.slides.add_slide(blank_slide_layout)
	#textbox (ignore, height where box is, size)
	txBox = slide.shapes.add_textbox(Inches(0), Inches(2.45), width, Inches(0.64))
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.alignment = PP_ALIGN.CENTER
	run = p.add_run()
	run.text = line
	font = run.font
	font.name = 'Source Han Serif K'
	font.size = Pt(28)
	font.color.rgb = RGBColor(255, 255, 255)

	txBox2 = slide.shapes.add_textbox(Inches(0), Inches(3.02), width, Inches(0.67))
	tf2 = txBox2.text_frame
	p2 = tf2.paragraphs[0]
	p2.alignment = PP_ALIGN.CENTER
	run2 = p2.add_run()
	run2.text = textfile2[count]
	count +=1
	font2 = run2.font
	font2.name = 'Source Han Serif K'
	font2.size = Pt(16.5)
	font2.color.rgb = RGBColor(255, 255, 255)

	txBox3 = slide.shapes.add_textbox(Inches(0), Inches(1.96), width, Inches(0.67))
	tf3 = txBox3.text_frame
	p3 = tf3.paragraphs[0]
	p3.alignment = PP_ALIGN.CENTER
	run3 = p3.add_run()
	run3.text = textfile3[count1]
	count1 +=1
	font3 = run3.font
	font3.name = 'Source Han Serif K'
	font3.size = Pt(25)
	font3.color.rgb = RGBColor(255, 255, 255)

prs.save('/Users/loser/Downloads/done_slides.pptx')

