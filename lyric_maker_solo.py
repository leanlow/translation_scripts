#! /usr/bin/env python
import csv, sys, numpy, math, pandas, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path

prs = Presentation()
prs.slide_height = Inches(5.63)
prs.slide_width = Inches(10)
blank_slide_layout = prs.slide_layouts[6]

width = Inches(7.28)
height = Inches(0.52)

textfile1 = open(sys.argv[1]).read().splitlines()
textfile2 = open(sys.argv[2]).read().splitlines()
textfile3 = open(sys.argv[3]).read().splitlines()

count = 0
count1 = 0

for line in textfile1:
	slide = prs.slides.add_slide(blank_slide_layout)
	#(left, top, width, height)
	txBox = slide.shapes.add_textbox(Inches(2.63), Inches(0.89), width, height)
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.alignment = PP_ALIGN.CENTER
	run = p.add_run()
	run.text = line
	font = run.font
	font.name = '210 Hayanbunpil Light'
	font.size = Pt(26)

	txBox2 = slide.shapes.add_textbox(Inches(2.63), Inches(2.45), width, height)
	tf2 = txBox2.text_frame
	p2 = tf2.paragraphs[0]
	p2.alignment = PP_ALIGN.CENTER
	run2 = p2.add_run()
	run2.text = textfile2[count]
	count +=1
	font2 = run2.font
	font2.name = '210 Hayanbunpil Light'
	font2.size = Pt(20)

	txBox3 = slide.shapes.add_textbox(Inches(2.63), Inches(4.01), width, height)
	tf3 = txBox3.text_frame
	p3 = tf3.paragraphs[0]
	p3.alignment = PP_ALIGN.CENTER
	run3 = p3.add_run()
	run3.text = textfile3[count1]
	count1 +=1
	font3 = run3.font
	font3.name = 'Throw My Hands Up in the Air Regular'
	font3.size = Pt(24)

prs.save('/Users/LeAnnLo/Downloads/finished_slides.pptx')

