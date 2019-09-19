#! /usr/bin/env python
import csv, sys, numpy, math, pandas, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path

prs = Presentation()
prs.slide_height = Inches(5.63)
prs.slide_width = Inches(10)
blank_slide_layout = prs.slide_layouts[6]

width = Inches(10)
#(left, top, width, height)
textfile1 = open(sys.argv[1]).read().splitlines()
textfile2 = open(sys.argv[2]).read().splitlines()

count = 0
count1 = 0

for line in textfile1:
	slide = prs.slides.add_slide(blank_slide_layout)
	#print("go")
	txBox = slide.shapes.add_textbox(Inches(0), Inches(3.36), width, Inches(0.64))
	tf = txBox.text_frame
	p = tf.paragraphs[0]
	p.alignment = PP_ALIGN.CENTER
	run = p.add_run()
	run.text = line
	font = run.font
	font.name = 'Throw My Hands Up in the Air Regular'
	font.size = Pt(28)

	txBox2 = slide.shapes.add_textbox(Inches(0), Inches(4.39), width, Inches(0.67))
	tf2 = txBox2.text_frame
	p2 = tf2.paragraphs[0]
	p2.alignment = PP_ALIGN.CENTER
	run2 = p2.add_run()
	run2.text = textfile2[count]
	count +=1
	font2 = run2.font
	font2.name = '210 Hayanbunpil Light'
	font2.size = Pt(25)

prs.save('/Users/loser/Downloads/done_slides.pptx')

